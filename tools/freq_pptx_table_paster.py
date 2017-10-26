from pptx import Presentation
from pptx.shapes.graphfrm import GraphicFrame
from collections import namedtuple

import pyperclip
import os
import sys

PPTX_PATH = 'tests/data/power_point/cycle_5.pptx'
# TODO: コマンドライン引数化
SLIDE_INDEX = 80
TABLE_INDEX = 0

RANK_NUMS = 10
START_ROW_INDEX = 2
END_ROW_INDEX = 12

# TODO: コマンドライン引数化
WORD_COLUMN_INDEX = 1
COUNT_COLUMN_INDEX = 2

FreqInfo = namedtuple('FreqInfo', ('rank', 'word', 'count'))

def extract_tables(slide):
    shapes = slide.shapes
    g_frames = [s for s in shapes if isinstance(s, GraphicFrame)]

    return [gf.table for gf in g_frames]

def get_target_table(prs):
    slides = prs.slides
    slide = slides[SLIDE_INDEX]
    tables = extract_tables(slide)
    table = tables[TABLE_INDEX]

    return table

def replace_paragraph_text_retaining_initial_formatting(paragraph, new_text):
    # 「https://github.com/scanny/python-pptx/issues/285」より

    p = paragraph._p  # the lxml element containing the `<a:p>` paragraph element
    # remove all but the first run
    for idx, run in enumerate(paragraph.runs):
        if idx == 0:
            continue
        p.remove(run._r)
    paragraph.runs[0].text = new_text

def replace_table_cell_text(cell, new_text):
    paragraph = cell.text_frame.paragraphs[0]
    replace_paragraph_text_retaining_initial_formatting(paragraph, new_text)

def create_single_freq_info(single_line):
    try:
        rank, word, count = single_line.split('\t')
    except ValueError:
        print('Clipboard data invalid.')
        sys.exit(1)

    return FreqInfo(rank, word, count)

def load_freq_infos():
    '''クリップボードから各単語出現数の情報を読み込む。

    :return:
    '''
    raw_clip = pyperclip.paste()
    lines = raw_clip.split('\n')

    return [create_single_freq_info(l) for l in lines]

def paste():
    '''クリップボードの「各単語出現数」をパワーポイントのテーブルに貼り付ける。

    クリップボードの形式は「(ランキング)\t(単語名)\t(出現数)」を改行文字で結合したもの。

    :return:
    '''
    prs = Presentation(PPTX_PATH)
    table = get_target_table(prs)

    cell = table.cell(2, 1)
    replace_table_cell_text(cell, 'aaa')
    # TODO: load_freq_infos等を使うようにする

    prs.save(PPTX_PATH[:-len('.pptx')] + '_modified.pptx')

if __name__ == '__main__':
    paste()


