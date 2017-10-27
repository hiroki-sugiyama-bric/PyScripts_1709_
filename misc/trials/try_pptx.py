from pptx import Presentation
from pptx.shapes.graphfrm import GraphicFrame

PPTX_PATH = 'data/power_point/cycle_5.pptx'
SLIDE_INDEX = 80

def extract_tables(slide):
    shapes = slide.shapes
    g_frames = [s for s in shapes if isinstance(s, GraphicFrame)]

    return [gf.table for gf in g_frames]

def replace_paragraph_text_retaining_initial_formatting(paragraph, new_text):
    # 「https://github.com/scanny/python-pptx/issues/285」より

    p = paragraph._p  # the lxml element containing the `<a:p>` paragraph element
    # remove all but the first run
    for idx, run in enumerate(paragraph.runs):
        if idx == 0:
            continue
        p.remove(run._r)
    paragraph.runs[0].text = new_text

def test_load():
    prs = Presentation(PPTX_PATH)
    slides = prs.slides
    slide = slides[SLIDE_INDEX]
    # shapes = slide.shapes
    #
    # for shape in shapes:
    #     print(shape)

    tables = extract_tables(slide)
    for table in tables:
        cell = table.cell(2, 1)
        original_text = cell.text_frame.text
        modified_text = original_text + '_modified'
        # cell.text_frame.text = modified_text
        paragraph = cell.text_frame.paragraphs[0]
        replace_paragraph_text_retaining_initial_formatting(paragraph, modified_text)


    prs.save(PPTX_PATH[:-len('.pptx')] + '_modified.pptx')


if __name__ == '__main__':
    test_load()


